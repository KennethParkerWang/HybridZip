from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
LAYOUT_PATH = ROOT / "layouts" / "page_001.layout.json"
OUTPUT_PATH = ROOT / "output" / "HybridZip_R2_target_architecture_editable.pptx"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_font(run, name: str, size: float, bold: bool, color: RGBColor) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", name)


def add_text(slide, text: str, box, size: float, color: RGBColor, bold=False,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02,
             name="text"):
    x, y, w, h = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 0.92
        run = p.add_run()
        run.text = line
        set_font(run, "Microsoft YaHei", size, bold, color)
    return shape


def add_box(slide, box, fill, line, width, dash=False, name="box"):
    x, y, w, h = box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.shadow.inherit = False
    return shape


def add_arrowhead(line) -> None:
    ln = line._element.spPr.ln
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    ln.append(tail)


def add_segment(slide, start, end, color, width, dash, arrow, name):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1])
    )
    line.name = name
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    line.shadow.inherit = False
    if arrow:
        add_arrowhead(line)
    return line


def add_route(slide, points, color, width, dash, name):
    for index in range(len(points) - 1):
        add_segment(
            slide,
            points[index],
            points[index + 1],
            color,
            width,
            dash,
            index == len(points) - 2,
            f"{name}-{index + 1}",
        )


def add_node(slide, node, colors):
    styles = {
        "current": (rgb("#FFFFFF"), colors["navy"], 2.2, False, colors["text"]),
        "focus": (colors["focus"], colors["navy"], 3.2, False, colors["text"]),
        "baseline": (rgb("#FFFFFF"), colors["gray"], 2.0, False, colors["text"]),
        "planned": (rgb("#FFFFFF"), colors["gray"], 2.0, True, colors["text"]),
        "current_dark": (colors["navy"], colors["navy"], 2.2, False, rgb("#FFFFFF")),
    }
    fill, line, width, dash, text_color = styles[node["style"]]
    x, y, w, h = node["box"]
    add_box(slide, node["box"], fill, line, width, dash, node["id"])
    title = node["title"]
    subtitle = node.get("subtitle")
    if subtitle:
        title_h = 0.38 if h < 1.2 else 0.42
        add_text(
            slide, title, [x + 0.05, y + 0.09, w - 0.10, title_h],
            node.get("title_size", 16), text_color, True, name=f"{node['id']}-title"
        )
        add_text(
            slide, subtitle, [x + 0.06, y + 0.48, w - 0.12, h - 0.53],
            node.get("subtitle_size", 12), text_color, False,
            name=f"{node['id']}-subtitle"
        )
    else:
        add_text(
            slide, title, [x + 0.05, y + 0.06, w - 0.10, h - 0.12],
            node.get("title_size", 16), text_color, True, name=f"{node['id']}-title"
        )


def build():
    layout = json.loads(LAYOUT_PATH.read_text(encoding="utf-8"))
    colors = {key: rgb(value) for key, value in layout["colors"].items()}
    prs = Presentation()
    prs.slide_width = Inches(layout["slide_size"][0])
    prs.slide_height = Inches(layout["slide_size"][1])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(layout["background"])

    add_text(
        slide, layout["title"], [0.30, 0.20, 5.8, 0.55],
        27, colors["navy"], True, PP_ALIGN.LEFT, name="slide-title"
    )

    # Legend: box + arrow samples make both border and path semantics explicit.
    add_box(slide, [9.15, 0.28, 0.28, 0.22], rgb("#FFFFFF"), colors["navy"], 1.8, False, "legend-current-box")
    add_segment(slide, [9.55, 0.39], [10.05, 0.39], colors["navy"], 2.2, False, True, "legend-current-arrow")
    add_text(slide, "当前已有基础", [10.10, 0.21, 1.05, 0.36], 12.5, colors["muted"], False, PP_ALIGN.LEFT)
    add_box(slide, [11.25, 0.28, 0.28, 0.22], rgb("#FFFFFF"), colors["gray"], 1.8, True, "legend-planned-box")
    add_segment(slide, [11.65, 0.39], [12.15, 0.39], colors["gray"], 2.0, True, True, "legend-planned-arrow")
    add_text(slide, "最终架构扩展", [12.20, 0.21, 1.00, 0.36], 12.5, colors["muted"], False, PP_ALIGN.LEFT)

    for idx, zone in enumerate(layout["zone_labels"], start=1):
        x, y, w, h = zone["box"]
        add_text(slide, zone["text"], zone["box"], zone.get("size", 11.5), colors["gray"], True, PP_ALIGN.LEFT, name=f"zone-{idx}")
        add_segment(slide, [x, y + h + 0.02], [x + w, y + h + 0.02], colors["gray"], 0.8, False, False, f"zone-rule-{idx}")

    # Add connectors before boxes so paths stay behind nodes.
    for idx, edge in enumerate(layout["edges"], start=1):
        style = edge["style"]
        add_route(
            slide,
            edge["points"],
            colors["navy"] if style == "current" else colors["gray"],
            2.4 if style == "current" else 2.0,
            style == "planned",
            f"edge-{idx}",
        )

    for node in layout["nodes"]:
        add_node(slide, node, colors)

    for idx, caption in enumerate(layout["captions"], start=1):
        add_text(slide, caption["text"], caption["box"], 10.5, colors["muted"], False, name=f"caption-{idx}")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    build()
